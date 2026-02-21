from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import ChartData as PPTChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pathlib import Path

from .schemas import PresentationSchema
from .export_utils import build_export_path

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    raise ImportError(
        "python-pptx is required for PowerPoint generation. "
        "Install with: pip install python-pptx"
    )



def generate_presentation(schema: PresentationSchema) -> str:

    prs = Presentation()

    layout_map = {
        "title": 0,
        "content": 1,
        "blank": 6,
    }

    for slide_data in schema.slides:

        layout_index = layout_map.get(slide_data.layout, 1)
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)

        # ----------------------------
        # Background Image
        # ----------------------------
        if slide_data.background_image:
            bg_path = Path(slide_data.background_image)
            if bg_path.exists():
                slide.background.fill.user_picture(str(bg_path))

        # ----------------------------
        # Title
        # ----------------------------
        if slide_data.title and slide.shapes.title:
            slide.shapes.title.text = slide_data.title

        # ----------------------------
        # Bullets
        # ----------------------------
        if slide_data.bullets and len(slide.placeholders) > 1:

            content_placeholder = slide.placeholders[1]
            text_frame = content_placeholder.text_frame
            text_frame.clear()

            for i, bullet in enumerate(slide_data.bullets):
                p = (
                    text_frame.paragraphs[0]
                    if i == 0
                    else text_frame.add_paragraph()
                )
                p.text = bullet.text
                p.level = max(0, min(bullet.level, 4))

        # ----------------------------
        # Table
        # ----------------------------
        if slide_data.table:

            rows = len(slide_data.table.rows) + 1
            cols = len(slide_data.table.headers)

            table = slide.shapes.add_table(
                rows,
                cols,
                Inches(1),
                Inches(2),
                Inches(8),
                Inches(4),
            ).table

            # Headers
            for col_index, header in enumerate(slide_data.table.headers):
                table.cell(0, col_index).text = header

            # Data
            for row_index, row in enumerate(slide_data.table.rows):
                for col_index, cell in enumerate(row):
                    table.cell(row_index + 1, col_index).text = str(cell)

        # ----------------------------
        # Chart
        # ----------------------------
        if slide_data.chart:

            categories = slide_data.chart.categories
            series_dict = slide_data.chart.series

            # Validate
            for name, values in series_dict.items():
                if len(values) != len(categories):
                    raise ValueError(
                        f"Series '{name}' length must match categories."
                    )

            chart_data = PPTChartData()
            chart_data.categories = categories

            for name, values in series_dict.items():
                chart_data.add_series(name, values)

            chart_type_map = {
                "line": XL_CHART_TYPE.LINE,
                "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            }

            chart_type = chart_type_map[slide_data.chart.chart_type]

            x, y = slide_data.chart.position
            width, height = slide_data.chart.size

            slide.shapes.add_chart(
                chart_type,
                Inches(x),
                Inches(y),
                Inches(width),
                Inches(height),
                chart_data,
            )

        # ----------------------------
        # Image
        # ----------------------------
        if slide_data.image_path:
            image_path = Path(slide_data.image_path)

            if not image_path.exists():
                raise FileNotFoundError(
                    f"Image not found: {slide_data.image_path}"
                )

            slide.shapes.add_picture(
                str(image_path),
                Inches(5),
                Inches(2),
                width=Inches(4),
            )

        # ----------------------------
        # Notes
        # ----------------------------
        if slide_data.notes:
            slide.notes_slide.notes_text_frame.text = slide_data.notes

        # ----------------------------
        # Footer
        # ----------------------------
        if schema.footer_text:
            for shape in slide.placeholders:
                if shape.is_placeholder:
                    continue
            slide.footer = schema.footer_text

    file_path = build_export_path(schema.file_name, "pptx")
    prs.save(file_path)

    return str(file_path)
