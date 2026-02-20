"""PowerPoint Tool - Generate .pptx slide decks (schema-first)."""
from __future__ import annotations

import os
from typing import Any, List

from fastmcp import FastMCP
from pydantic import BaseModel, Field

from aden_tools.tools.office_skills_pack.contracts import ArtifactError, ArtifactResult
from ..file_system_toolkits.security import get_secure_path


class SlideSpec(BaseModel):
    title: str = Field(..., min_length=1)
    bullets: List[str] = Field(default_factory=list)
    image_paths: List[str] = Field(default_factory=list, description="Paths relative to session sandbox")

    charts: List[str] = Field(
        default_factory=list,
        description="Chart PNG paths relative to session sandbox"
    )


class DeckSpec(BaseModel):
    title: str = Field(..., min_length=1)
    slides: List[SlideSpec] = Field(..., min_length=1)


def register_tools(mcp: FastMCP) -> None:
    @mcp.tool()
    def powerpoint_generate(
        path: str,
        deck: dict[str, Any],
        workspace_id: str,
        agent_id: str,
        session_id: str,
        strict: bool = True,
    ) -> dict:
        """
        Generate a PowerPoint (.pptx) from a strict schema and save it to the session sandbox.

        Args:
            path: Output path (relative to session sandbox), must end with .pptx
            deck: Dict matching DeckSpec schema
            workspace_id/agent_id/session_id: sandbox identifiers

        Returns:
            ArtifactResult dict with success, output_path, and metadata
        """
        if not path.lower().endswith(".pptx"):
            return ArtifactResult(
                success=False,
                error=ArtifactError(
                    code="INVALID_SCHEMA",
                    message="path must end with .pptx",
                ),
            ).model_dump()

        try:
            spec = DeckSpec.model_validate(deck)
        except Exception as e:
            return ArtifactResult(
                success=False,
                error=ArtifactError(
                    code="INVALID_SCHEMA",
                    message="Invalid deck schema",
                    details=str(e),
                ),
            ).model_dump()

        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            return ArtifactResult(
                success=False,
                error=ArtifactError(
                    code="DEP_MISSING",
                    message="python-pptx not installed. Install with: pip install -e '.[powerpoint]' (from tools/)",
                ),
            ).model_dump()

        try:
            out_path = get_secure_path(path, workspace_id, agent_id, session_id)
            os.makedirs(os.path.dirname(out_path), exist_ok=True)

            prs = Presentation()

            # Title slide
            title_layout = prs.slide_layouts[0] if len(prs.slide_layouts) > 0 else prs.slide_layouts[5]
            slide0 = prs.slides.add_slide(title_layout)
            if slide0.shapes.title:
                slide0.shapes.title.text = spec.title

            # Content slides
            content_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[5]

            for s in spec.slides:
                sl = prs.slides.add_slide(content_layout)
                if sl.shapes.title:
                    sl.shapes.title.text = s.title

                # Find a body text frame (best-effort)
                body = None
                for shape in sl.shapes:
                    if shape.has_text_frame and shape != sl.shapes.title:
                        body = shape.text_frame
                        break

                if body and s.bullets:
                    body.clear()
                    p0 = body.paragraphs[0]
                    p0.text = s.bullets[0]
                    p0.level = 0
                    for b in s.bullets[1:]:
                        p = body.add_paragraph()
                        p.text = b
                        p.level = 0

                all_imgs = list(s.image_paths) + list(s.charts)
                positions = [(1.0, 4.5), (3.7, 4.5), (6.4, 4.5)]  # left-to-right row

                for idx, img_rel in enumerate(all_imgs[:3]):
                    img_abs = get_secure_path(img_rel, workspace_id, agent_id, session_id)
                    if not os.path.exists(img_abs):
                        if strict:
                            return ArtifactResult(
                                success=False,
                                error=ArtifactError(
                                    code="INVALID_PATH",
                                    message=f"image not found: {img_rel}",
                                ),
                            ).model_dump()
                        continue
                    x, y = positions[idx]
                    sl.shapes.add_picture(img_abs, Inches(x), Inches(y), width=Inches(2.5))

            prs.save(out_path)

            return ArtifactResult(
                success=True,
                output_path=str(out_path),
                metadata={"slides": len(prs.slides)},
            ).model_dump()

        except Exception as e:
            return ArtifactResult(
                success=False,
                error=ArtifactError(
                    code="INTERNAL_ERROR",
                    message="Failed to generate PowerPoint",
                    details=str(e),
                ),
            ).model_dump()
