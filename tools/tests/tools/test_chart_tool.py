import importlib.util
from pathlib import Path
from unittest.mock import patch

import pytest
from fastmcp import FastMCP

from aden_tools.tools.chart_tool.chart_tool import register_tools as register_chart
from aden_tools.tools.powerpoint_tool.powerpoint_tool import register_tools as register_ppt
from aden_tools.tools.file_system_toolkits.security import get_secure_path

mpl_available = importlib.util.find_spec("matplotlib") is not None
pytestmark = pytest.mark.skipif(not mpl_available, reason="matplotlib not installed")


def test_chart_render_png_and_embed_in_pptx(tmp_path: Path):
    with patch("aden_tools.tools.file_system_toolkits.security.WORKSPACES_DIR", str(tmp_path)):
        mcp = FastMCP("chart-test")
        register_chart(mcp)
        register_ppt(mcp)

        chart_png = mcp._tool_manager._tools["chart_render_png"].fn
        pptx = mcp._tool_manager._tools["powerpoint_generate"].fn

        r = chart_png(
            path="out/chart.png",
            chart={
                "title": "Test Chart",
                "x_label": "x",
                "y_label": "y",
                "series": [{"name": "s1", "x": [1, 2, 3], "y": [1, 4, 9]}],
            },
            workspace_id="w",
            agent_id="a",
            session_id="s",
        )
        assert r["success"] is True, r
        assert len(r.get("metadata", {}).get("sha256", "")) == 64

        chart_abs = get_secure_path("out/chart.png", "w", "a", "s")
        assert Path(chart_abs).exists()
        assert Path(chart_abs).stat().st_size > 0

        r2 = pptx(
            path="out/deck.pptx",
            deck={
                "title": "Deck",
                "slides": [
                    {"title": "Chart", "bullets": ["see chart"], "image_paths": ["out/chart.png"], "charts": []},
                ],
            },
            workspace_id="w",
            agent_id="a",
            session_id="s",
        )
        assert r2["success"] is True, r2

        ppt_abs = get_secure_path("out/deck.pptx", "w", "a", "s")
        assert Path(ppt_abs).exists()

        from pptx import Presentation
        prs = Presentation(str(ppt_abs))

        # Check: at least one picture shape exists somewhere
        has_picture = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE == 13
                    has_picture = True
        assert has_picture is True
